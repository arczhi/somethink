"""
文件扫描模块
负责扫描指定目录下的文件，提取元数据和文本内容
"""

import os
from pathlib import Path
from typing import List, Dict, Generator
from datetime import datetime
import mimetypes


class FileScanner:
    """文件扫描器"""
    
    # 支持的文件类型
    SUPPORTED_EXTENSIONS = {
        'document': ['.txt', '.md', '.pdf', '.docx', '.pptx', '.doc', '.rtf'],
        'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp'],
        'audio': ['.mp3', '.flac', '.wav', '.m4a', '.aac', '.ogg'],
        'video': ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm']
    }
    
    def __init__(self, exclude_patterns: List[str] = None):
        self.exclude_patterns = exclude_patterns or []
        
    def should_exclude(self, path: Path) -> bool:
        """判断路径是否应该被排除"""
        for pattern in self.exclude_patterns:
            if pattern in path.parts:
                return True
        return False
    
    def get_file_type(self, file_path: Path) -> str:
        """获取文件类型"""
        ext = file_path.suffix.lower()
        
        for file_type, extensions in self.SUPPORTED_EXTENSIONS.items():
            if ext in extensions:
                return file_type
        
        return 'unknown'
    
    def is_supported(self, file_path: Path) -> bool:
        """判断文件是否被支持"""
        return self.get_file_type(file_path) != 'unknown'
    
    def scan_directory(self, root_path: str) -> Generator[Dict, None, None]:
        """
        扫描目录，返回文件信息生成器
        
        Args:
            root_path: 根目录路径
            
        Yields:
            文件信息字典
        """
        root = Path(root_path)
        
        if not root.exists() or not root.is_dir():
            return
        
        for file_path in root.rglob('*'):
            # 跳过目录和排除的路径
            if file_path.is_dir() or self.should_exclude(file_path):
                continue
            
            # 只处理支持的文件类型
            if not self.is_supported(file_path):
                continue
            
            try:
                file_info = self.get_file_info(file_path)
                yield file_info
            except Exception as e:
                print(f"处理文件 {file_path} 时出错: {e}")
                continue
    
    def get_file_info(self, file_path: Path) -> Dict:
        """
        获取文件基本信息
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件信息字典
        """
        stat = file_path.stat()
        
        return {
            'path': str(file_path.absolute()),
            'filename': file_path.name,
            'extension': file_path.suffix.lower(),
            'file_type': self.get_file_type(file_path),
            'size': stat.st_size,
            'modified_time': datetime.fromtimestamp(stat.st_mtime),
            'created_time': datetime.fromtimestamp(stat.st_ctime),
        }
    
    def extract_text_content(self, file_path: Path) -> str:
        """
        提取文件文本内容
        
        Args:
            file_path: 文件路径
            
        Returns:
            提取的文本内容
        """
        file_type = self.get_file_type(file_path)
        
        if file_type == 'document':
            return self._extract_document_text(file_path)
        elif file_type == 'image':
            return self._extract_image_text(file_path)
        elif file_type == 'audio':
            return self._extract_audio_metadata(file_path)
        elif file_type == 'video':
            return self._extract_video_metadata(file_path)
        
        return ""
    
    def _extract_document_text(self, file_path: Path) -> str:
        """提取文档文本"""
        ext = file_path.suffix.lower()
        
        try:
            if ext == '.txt' or ext == '.md':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            elif ext == '.pdf':
                return self._extract_pdf_text(file_path)
            
            elif ext in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)
            
            elif ext == '.pptx':
                return self._extract_pptx_text(file_path)
        
        except Exception as e:
            print(f"提取文档 {file_path} 文本失败: {e}")
        
        return ""
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """提取PDF文本"""
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(str(file_path))
            text_parts = []
            
            for page in reader.pages[:10]:  # 限制前10页，避免过大
                text_parts.append(page.extract_text())
            
            return "\n".join(text_parts)
        except ImportError:
            print("PyPDF2未安装，跳过PDF文本提取")
            return ""
        except Exception as e:
            print(f"PDF提取失败: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: Path) -> str:
        """提取DOCX文本"""
        try:
            from docx import Document
            
            doc = Document(str(file_path))
            text_parts = [para.text for para in doc.paragraphs]
            
            return "\n".join(text_parts)
        except ImportError:
            print("python-docx未安装，跳过DOCX文本提取")
            return ""
        except Exception as e:
            print(f"DOCX提取失败: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: Path) -> str:
        """提取PPTX文本"""
        try:
            from pptx import Presentation
            
            prs = Presentation(str(file_path))
            text_parts = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
        except ImportError:
            print("python-pptx未安装，跳过PPTX文本提取")
            return ""
        except Exception as e:
            print(f"PPTX提取失败: {e}")
            return ""
    
    def _extract_image_text(self, file_path: Path) -> str:
        """提取图片元数据（未来可添加OCR）"""
        try:
            from PIL import Image
            from PIL.ExifTags import TAGS
            
            image = Image.open(str(file_path))
            exif_data = image.getexif()
            
            metadata_parts = []
            
            if exif_data:
                for tag_id, value in exif_data.items():
                    tag = TAGS.get(tag_id, tag_id)
                    metadata_parts.append(f"{tag}: {value}")
            
            # 添加基本信息
            metadata_parts.append(f"Size: {image.size}")
            metadata_parts.append(f"Mode: {image.mode}")
            
            return " | ".join(metadata_parts)
        
        except Exception as e:
            return ""
    
    def _extract_audio_metadata(self, file_path: Path) -> str:
        """提取音频元数据"""
        try:
            from mutagen import File
            
            audio = File(str(file_path))
            
            if audio is None:
                return ""
            
            metadata_parts = []
            
            # 常见标签
            tags = ['title', 'artist', 'album', 'genre', 'date']
            
            for tag in tags:
                if tag in audio and audio[tag]:
                    value = audio[tag][0] if isinstance(audio[tag], list) else audio[tag]
                    metadata_parts.append(f"{tag}: {value}")
            
            return " | ".join(metadata_parts)
        
        except ImportError:
            print("mutagen未安装，跳过音频元数据提取")
            return ""
        except Exception as e:
            return ""
    
    def _extract_video_metadata(self, file_path: Path) -> str:
        """提取视频元数据（简化版）"""
        # 视频元数据提取较复杂，这里仅返回文件名作为内容
        return file_path.stem


if __name__ == "__main__":
    # 测试代码
    scanner = FileScanner(exclude_patterns=['.git', 'node_modules'])
    
    test_path = Path.home() / "Documents"
    
    count = 0
    for file_info in scanner.scan_directory(str(test_path)):
        print(f"发现文件: {file_info['filename']} ({file_info['file_type']})")
        count += 1
        
        if count >= 10:  # 仅测试前10个文件
            break
